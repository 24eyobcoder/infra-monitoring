#!/usr/bin/env python3
"""Generate the InfraWatch prototype demo deck (PowerPoint)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (matches the dashboard) ----
BG      = RGBColor(0x0A, 0x0E, 0x14)
PANEL   = RGBColor(0x12, 0x18, 0x22)
PANEL2  = RGBColor(0x1B, 0x22, 0x2E)
LINE    = RGBColor(0x2A, 0x35, 0x43)
TEXT    = RGBColor(0xE4, 0xE9, 0xF0)
DIM     = RGBColor(0x9A, 0xA5, 0xB4)
BRAND   = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)
MONO    = "Consolas"
SANS    = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=4, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.line_spacing = line_spacing
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
    return tb


def kicker(s, label):
    box(s, 0.0, 0.0, 0.18, SH.inches, fill=BRAND)  # left accent rail
    text(s, 0.7, 0.55, 11, 0.4,
         [[("INFRAWATCH", 12, BRAND, True, MONO), ("  ·  " + label, 12, DIM, False, MONO)]])


def title(s, t, sub=None):
    kicker(s, sub or "")
    text(s, 0.7, 1.0, 12, 1.0, [[(t, 32, TEXT, True, SANS)]])


def bullets(s, items, x=0.9, y=2.1, w=11.5, size=17, gap=10):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            runs.append([("●  ", 15, BRAND, True, SANS), (head, size, TEXT, True, SANS),
                         ("  —  " + body, size, DIM, False, SANS)])
        else:
            runs.append([("●  ", 15, BRAND, True, SANS), (it, size, TEXT, False, SANS)])
    text(s, x, y, w, 5, runs, space=gap, line_spacing=1.1)


def chip(s, x, y, w, label, color):
    h = 0.55
    b = box(s, x, y, w, h, fill=PANEL, line=color, line_w=1.25, round_=True)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = color; r.font.name = MONO
    return b


# ============================================================ 1. TITLE
s = slide()
box(s, 0, 0, SW.inches, 0.10, fill=BRAND)
# radar mark
box(s, 0.9, 2.35, 0.7, 0.7, fill=PANEL2, line=BRAND, line_w=1.5, round_=True)
text(s, 0.9, 2.5, 0.7, 0.4, [[("◎", 26, BRAND, True, SANS)]], align=PP_ALIGN.CENTER)
text(s, 1.8, 2.25, 11, 1.4, [[("InfraWatch", 54, TEXT, True, SANS)]])
text(s, 1.82, 3.45, 11, 0.6,
     [[("AI-Augmented Infrastructure Monitoring", 22, BRAND, False, SANS)]])
text(s, 1.82, 4.15, 11, 0.5,
     [[("Live device health · AI-summarized alerts · Telegram push · Reporting", 15, DIM, False, SANS)]])
text(s, 1.82, 5.7, 11, 0.5,
     [[("PROTOTYPE DEMO", 13, AMBER, True, MONO), ("    eTech", 13, DIM, False, MONO)]])

# ============================================================ 2. PROBLEM
s = slide(); title(s, "The problem", "Why we built it")
bullets(s, [
    ("Raw monitoring alerts are cryptic", "Munin emits lines like \"temp1 is 89.00 (outside range [:70.0])\" — hard to read under pressure."),
    ("Alerts are scattered", "buried in emails/logs; no single place to see what's on fire right now."),
    ("No fast human context", "on-call staff waste time decoding metrics instead of acting."),
    ("No history for review", "no record of what happened, when, or how often."),
], y=2.2, gap=14)

# ============================================================ 3. SOLUTION
s = slide(); title(s, "The solution", "What InfraWatch does")
bullets(s, [
    ("Watches infrastructure", "pings hosts, checks HTTP/ports/SNMP/UPS, shows live up / degraded / down status."),
    ("Understands Munin alerts", "receives threshold breaches via a webhook from the existing Munin setup."),
    ("Explains them with AI", "Google Gemini rewrites each alert into a plain-English summary + suggested action."),
    ("Notifies instantly", "pushes the summary to a Telegram channel for the on-call team."),
    ("Remembers everything", "stores every alert with a timestamp in SQLite for a dedicated report page."),
], y=1.95, gap=12)

# ============================================================ 4. ARCHITECTURE
s = slide(); title(s, "Architecture & data flow", "How it fits together")
# pipeline boxes
y0 = 2.35; bw = 2.05; bh = 1.0; gap = 0.35; x = 0.7
steps = [("Munin", "threshold\nbreach", AMBER), ("Webhook", "munin-webhook.py", DIM),
         ("Express API", "/api/munin-alert", BRAND), ("Gemini AI", "summarize\n+ fallback", GREEN),
         ("SQLite", "store w/ time", BRAND)]
centers = []
for i, (h, sub, col) in enumerate(steps):
    bx = x + i * (bw + gap)
    box(s, bx, y0, bw, bh, fill=PANEL, line=col, line_w=1.5, round_=True)
    text(s, bx, y0 + 0.14, bw, 0.4, [[(h, 15, col, True, SANS)]], align=PP_ALIGN.CENTER)
    text(s, bx, y0 + 0.52, bw, 0.4, [[(sub, 10.5, DIM, False, MONO)]], align=PP_ALIGN.CENTER)
    centers.append((bx + bw, y0 + bh / 2))
    if i < len(steps) - 1:
        text(s, bx + bw, y0 + 0.28, gap, 0.4, [[("→", 20, BRAND, True, SANS)]], align=PP_ALIGN.CENTER)
# outputs row
oy = 4.2
box(s, 8.05, oy, bw, 0.9, fill=PANEL, line=GREEN, line_w=1.5, round_=True)
text(s, 8.05, oy + 0.1, bw, 0.4, [[("Dashboard", 14, GREEN, True, SANS)]], align=PP_ALIGN.CENTER)
text(s, 8.05, oy + 0.46, bw, 0.4, [[("live + top criticals", 10, DIM, False, MONO)]], align=PP_ALIGN.CENTER)
box(s, 10.45, oy, bw, 0.9, fill=PANEL, line=BRAND, line_w=1.5, round_=True)
text(s, 10.45, oy + 0.1, bw, 0.4, [[("Telegram", 14, BRAND, True, SANS)]], align=PP_ALIGN.CENTER)
text(s, 10.45, oy + 0.46, bw, 0.4, [[("channel push", 10, DIM, False, MONO)]], align=PP_ALIGN.CENTER)
text(s, 9.0, 5.25, 3.2, 0.4, [[("SQLite feeds both →", 11, DIM, False, MONO)]], align=PP_ALIGN.CENTER)
# second flow: monitor loop
text(s, 0.7, 5.75, 12, 0.5,
     [[("Parallel: ", 14, AMBER, True, SANS),
       ("a 30s monitor loop checks every device (ping / HTTP / port / SNMP / UPS) → ", 14, TEXT, False, SANS),
       ("/api/status", 13, BRAND, False, MONO),
       (" → live dashboard.", 14, TEXT, False, SANS)]])

# ============================================================ 5. TECH STACK
s = slide(); title(s, "Tech stack", "Built with")
rows = [
    ("Backend", "Node.js · Express · nodemon", BRAND),
    ("AI layer", "Google Gemini (gemini-2.5-flash) via @google/genai", GREEN),
    ("Storage", "SQLite (better-sqlite3) — alerts with timestamps", BRAND),
    ("Alerting source", "Munin (munin-node) + Python webhook contact script", AMBER),
    ("Notifications", "Telegram Bot API (IPv4-forced HTTPS)", BRAND),
    ("Checks", "ICMP ping · HTTP · TCP port · SNMP · UPS battery", GREEN),
    ("Frontend", "Vanilla HTML / CSS / JS — dashboard + report page", BRAND),
]
y = 2.05
for label, val, col in rows:
    box(s, 0.9, y, 2.6, 0.62, fill=PANEL2, line=None, round_=True)
    text(s, 0.9, y + 0.13, 2.6, 0.4, [[(label, 13, col, True, MONO)]], align=PP_ALIGN.CENTER)
    text(s, 3.75, y + 0.12, 8.6, 0.45, [[(val, 15, TEXT, False, SANS)]])
    y += 0.72

# ============================================================ 6. DASHBOARD
s = slide(); title(s, "Live dashboard", "Feature 1")
bullets(s, [
    ("Real-time status grid", "every device shown as Online / Degraded / Offline, grouped, with uptime bars."),
    ("Summary tiles & banner", "instant count of total / up / degraded / down."),
    ("Critical Issues panel", "shows the 3 most recent CRITICAL alerts with their AI summary."),
    ("Filter & search", "issues-only view, search by name or IP; light/dark theme."),
    ("Auto-refresh", "polls every 10 seconds — no reload needed."),
], y=2.1, gap=12)

# ============================================================ 7. AI SUMMARIES
s = slide(); title(s, "AI alert summaries", "Feature 2 — the core idea")
text(s, 0.9, 1.95, 5.7, 0.4, [[("RAW MUNIN (in)", 12, DIM, True, MONO)]])
box(s, 0.9, 2.35, 5.7, 1.7, fill=PANEL, line=RED, line_w=1.25, round_=True)
text(s, 1.1, 2.55, 5.3, 1.4,
     [[("localhost.localdomain :: Temperatures", 12, TEXT, False, MONO)],
      [("CRITICALs: temp1 is 89.00", 12, RED, False, MONO)],
      [("(outside range [:70.0]).", 12, RED, False, MONO)]], line_spacing=1.2)
text(s, 6.95, 1.95, 5.6, 0.4, [[("AI SUMMARY (out)", 12, DIM, True, MONO)]])
box(s, 6.95, 2.35, 5.6, 1.7, fill=PANEL, line=GREEN, line_w=1.25, round_=True)
text(s, 7.15, 2.5, 5.2, 1.5,
     [[("The temperature on localhost.localdomain is "
        "critically high at 89°C, exceeding the 70°C "
        "limit. Please check the server's cooling system.", 14, TEXT, False, SANS)]], line_spacing=1.15)
bullets(s, [
    ("Prompted as an ops assistant", "plain English, names the host & metric, ends with a next step."),
    ("Resilient", "1 retry on transient errors, fails fast on quota/4xx, then a numbers-based fallback."),
    ("Never blocks", "if AI is down, the alert still shows the real readings."),
], y=4.45, gap=9, size=15)

# ============================================================ 8. MUNIN
s = slide(); title(s, "Munin integration", "Feature 3")
bullets(s, [
    ("Uses existing Munin setup", "no new agents — hooks into Munin's own threshold alerting."),
    ("Contact webhook script", "munin.conf runs munin-webhook.py on every warning/critical."),
    ("Parses the alert", "extracts host, title, severity and the real problem lines (drops healthy 'OK' readings)."),
    ("POSTs clean JSON", "to /api/munin-alert — the single entry point for the AI + storage pipeline."),
], y=2.15, gap=14)

# ============================================================ 9. TELEGRAM
s = slide(); title(s, "Telegram notifications", "Feature 4")
bullets(s, [
    ("Instant channel push", "every alert's AI summary (or fallback) posts to a Telegram channel."),
    ("Severity at a glance", "color-coded emoji, host, metric and timestamp in each message."),
    ("Fire-and-forget", "never blocks or breaks the webhook; logs its own errors."),
    ("Hardened delivery", "forces IPv4 — Telegram's IPv6 path was unreachable on the test network."),
], y=2.15, gap=14)

# ============================================================ 10. REPORTING
s = slide(); title(s, "Persistence & reporting", "Feature 5")
bullets(s, [
    ("Every alert is stored", "SQLite row with severity, host, metric, AI summary and ISO timestamp."),
    ("Dedicated report page", "/report — totals, by-severity, by-host, daily trend, full alert table."),
    ("Time windows", "24h / 7d / 30d / all-time presets; CSV export for sharing."),
    ("Dismiss ≠ delete", "dismissing clears the dashboard but keeps the row for the report."),
], y=2.15, gap=14)

# ============================================================ 11. RESILIENCE
s = slide(); title(s, "Engineering for resilience", "Built to degrade gracefully")
bullets(s, [
    ("AI down? Still useful", "fallback message shows the actual breached readings (the numbers)."),
    ("Smart retries", "retry transient network/5xx errors; fail fast on bad-key/quota (no wasted waits)."),
    ("Free-tier aware", "handles Gemini's daily quota — falls back instead of erroring."),
    ("Network-hardened", "IPv4-forced Telegram calls after diagnosing a broken IPv6 route."),
    ("Non-blocking pipeline", "storage and notification never delay the Munin webhook response."),
], y=1.95, gap=11)

# ============================================================ 12. DEMO SCRIPT
s = slide(); title(s, "Live demo", "What to show")
items = [
    ("Open the dashboard", "live device grid, status tiles, theme toggle."),
    ("Trigger a Munin alert", "force a check (sudo -u munin munin-cron) or POST a sample payload."),
    ("Watch it flow", "AI summary appears in the Critical Issues panel within seconds."),
    ("Check Telegram", "the same summary lands in the channel."),
    ("Open /report", "totals, trend, full history; export CSV."),
]
y = 2.1
for i, (h, b) in enumerate(items, 1):
    box(s, 0.9, y, 0.55, 0.55, fill=PANEL2, line=BRAND, line_w=1.25, round_=True)
    text(s, 0.9, y + 0.09, 0.55, 0.4, [[(str(i), 16, BRAND, True, MONO)]], align=PP_ALIGN.CENTER)
    text(s, 1.7, y + 0.02, 10.8, 0.6,
         [[(h, 17, TEXT, True, SANS), ("  —  " + b, 15, DIM, False, SANS)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78

# ============================================================ 13. ROADMAP
s = slide(); title(s, "Roadmap", "Prototype → production")
bullets(s, [
    ("Deploy to a VPS", "always-on host for the monitor loop + persistent SQLite (serverless can't run either)."),
    ("Scale the data layer", "move to Postgres / hosted DB as alert volume grows."),
    ("Higher AI throughput", "enable Gemini billing to lift the free-tier daily cap."),
    ("More notify channels", "email / Slack / on-call escalation alongside Telegram."),
    ("Auth & multi-tenant", "logins, per-team channels, role-based views."),
    ("Alert correlation", "group related alerts; let the AI spot patterns across hosts."),
], y=1.95, gap=11)

# ============================================================ 14. CLOSE
s = slide()
box(s, 0, 0, SW.inches, 0.10, fill=BRAND)
text(s, 0.9, 2.7, 11.5, 1.2, [[("Thank you", 46, TEXT, True, SANS)]])
text(s, 0.92, 3.85, 11.5, 0.6,
     [[("InfraWatch — turning raw infra alerts into clear, actionable, AI-summarized signals.", 18, BRAND, False, SANS)]])
text(s, 0.92, 4.7, 11.5, 0.5, [[("Questions & live demo", 15, DIM, False, MONO)]])

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "InfraWatch-Demo.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
